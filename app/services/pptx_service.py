from __future__ import annotations

import logging
from pathlib import Path
from typing import Dict

from pptx import Presentation
from io import BytesIO

from pptx.util import Inches, Pt

from app.config.slide_bindings import SLIDE_BINDINGS
from app.config.image_bindings import IMAGE_BINDINGS

logger = logging.getLogger(__name__)


def _find_shape(slide, shape_name: str):
    for shape in slide.shapes:
        if shape.name == shape_name:
            return shape
    return None


def _apply_text(shape, text: str, font_size: int | None = None, font_bold: bool | None = None) -> None:
    if not shape.has_text_frame:
        raise ValueError(f"Shape {shape.name} has no text frame")
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    if font_bold is not None:
        run.font.bold = font_bold
    if font_size is not None:
        run.font.size = Pt(font_size)
        return
    length = len(text)
    if length > 400:
        run.font.size = Pt(10)
    elif length > 250:
        run.font.size = Pt(12)
    else:
        run.font.size = Pt(14)


def _apply_image(slide, image_bytes: bytes, placement: Dict[str, float]) -> None:
    left = Inches(placement["left"])
    top = Inches(placement["top"])
    width = Inches(placement["width"])
    height = Inches(placement["height"])
    slide.shapes.add_picture(BytesIO(image_bytes), left, top, width=width, height=height)


def render_pptx(
    slide_draft: Dict[str, Dict[str, str]],
    template_path: str,
    output_path: str,
    image_bytes: bytes | None = None,
) -> None:
    logger.info("PPTX render start")
    prs = Presentation(template_path)

    for slide_key, bindings in SLIDE_BINDINGS.items():
        for field, target in bindings.items():
            slide_index = target["slide_index"]
            shape_name = target["shape_name"]
            try:
                slide = prs.slides[slide_index]
            except IndexError as exc:
                raise ValueError(f"Slide index {slide_index} not found") from exc

            shape = _find_shape(slide, shape_name)
            if not shape:
                raise ValueError(f"Shape not found: {shape_name}")

            text = slide_draft.get(slide_key, {}).get(field, "")
            font_size = target.get("font_size")
            font_bold = target.get("font_bold")
            _apply_text(shape, text, font_size=font_size, font_bold=font_bold)

    if image_bytes:
        for slide_key, placement in IMAGE_BINDINGS.items():
            slide_index = placement.get("slide_index")
            if slide_index is None:
                continue
            slide = prs.slides[slide_index]
            _apply_image(slide, image_bytes, placement)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    logger.info("PPTX render done")
